"""
生成智能审图汇报演示文档 (PPT)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# 颜色定义
BLUE       = RGBColor(0x00, 0x62, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
GREEN      = RGBColor(0x05, 0x96, 0x69)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG    = RGBColor(0xF7, 0xF8, 0xFA)
GRAY_TEXT  = RGBColor(0x4B, 0x55, 0x63)
MUTED      = RGBColor(0x9C, 0xA3, 0xAF)
BORDER     = RGBColor(0xE5, 0xE7, 0xEB)
ORANGE     = RGBColor(0xD9, 0x77, 0x06)
ORANGE_BG  = RGBColor(0xFE, 0xF3, 0xC7)
GREEN_BG   = RGBColor(0xDC, 0xFC, 0xE7)
BLUE_BG    = RGBColor(0xEF, 0xF6, 0xFF)
RED_DARK   = RGBColor(0xC2, 0x41, 0x0C)
RED_BG     = RGBColor(0xFF, 0xF7, 0xED)

W = Inches(13.33)   # 16:9 宽度
H = Inches(7.5)     # 16:9 高度

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]  # 完全空白

# ─────────────────── 辅助函数 ───────────────────

def add_rect(slide, x, y, w, h, fill_color=None, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color and border_pt > 0:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Arial"
    return txBox

def slide_header(slide, label, label_color, title_text):
    """通用页面标题区"""
    # 顶部装饰条
    add_rect(slide, 0, 0, 13.33, 0.08, label_color)
    # 标签badge
    add_rect(slide, 0.7, 0.25, 1.2, 0.38, label_color)
    add_text(slide, label, 0.75, 0.27, 1.1, 0.34, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 标题
    add_text(slide, title_text, 2.05, 0.18, 10.5, 0.6, 32, bold=True, color=DARK)
    # 分割线
    add_rect(slide, 0.7, 0.8, 11.93, 0.02, BORDER)


# ═══════════════════════════════════════════════
# 幻灯片 1 — 封面
# ═══════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)

# 左侧蓝色背景
add_rect(s1, 0, 0, 4.7, 7.5, BLUE)
# 右侧灰色背景
add_rect(s1, 4.7, 0, 8.63, 7.5, GRAY_BG)

# 左侧内容
add_rect(s1, 0.7, 1.0, 2.0, 0.38, WHITE)
add_text(s1, "AI 智能审图系统", 0.75, 1.02, 1.9, 0.34, 13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s1, "智能图纸\n审查平台", 0.7, 1.6, 3.7, 1.6, 48, bold=True, color=WHITE)
add_text(s1, "端到端自动化审图 · 智能解析 · 合规报告生成", 0.7, 3.4, 3.6, 0.7, 17, color=WHITE)
add_rect(s1, 0.7, 4.2, 0.55, 0.06, WHITE)
add_text(s1, "演示汇报文档 · 2025", 0.7, 4.35, 3.5, 0.35, 14, color=WHITE, italic=True)
add_text(s1, "版本 v2.0 | 审图平台团队", 0.7, 4.75, 3.5, 0.35, 13, color=RGBColor(0xBF,0xDB,0xFF))

# 右侧 KPI 卡片
add_text(s1, "系统能力概览", 5.1, 0.6, 7.5, 0.5, 24, bold=True, color=DARK)
add_text(s1, "本平台基于人工智能与规则引擎，实现建筑图纸的全自动化审查", 5.1, 1.15, 7.8, 0.55, 16, color=GRAY_TEXT)

# KPI 3栏
kpi_data = [("98%","审图准确率"),("<30s","单图处理时间"),("200+","规范条文覆盖")]
for i,(num,lbl) in enumerate(kpi_data):
    xk = 5.1 + i*2.6
    add_rect(s1, xk, 1.85, 2.3, 1.3, WHITE)
    add_text(s1, num, xk+0.15, 1.9, 2.0, 0.75, 36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s1, lbl, xk+0.15, 2.7, 2.0, 0.35, 14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# 功能卡片 3栏
feat_data = [("智能解析","自动识别图层、标注、尺寸"),("规则引擎","多维度合规性逐条比对"),("报告生成","结构化报告一键导出")]
for i,(ft,fd) in enumerate(feat_data):
    xf = 5.1 + i*2.6
    add_rect(s1, xf, 3.35, 2.3, 1.3, WHITE)
    add_text(s1, f"⬡ {ft}", xf+0.15, 3.42, 2.0, 0.4, 16, bold=True, color=DARK)
    add_text(s1, fd, xf+0.15, 3.85, 2.0, 0.7, 13, color=GRAY_TEXT)

add_text(s1, "SparkFlow 智能审图平台", 5.1, 7.1, 8.0, 0.3, 12, color=MUTED, italic=True)


# ═══════════════════════════════════════════════
# 幻灯片 2 — 架构功能流程图
# ═══════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, 13.33, 7.5, WHITE)
slide_header(s2, "系统架构", BLUE, "架构功能流程图")

cols = [
    ("① 输入层",   BLUE,  BLUE_BG,  [("DWG / CAD 文件","支持 AutoCAD 2000-2024"),("PDF 图纸","矢量 & 扫描版均支持"),("图片格式","PNG / JPG / TIFF"),("规范数据库","国标 / 地标 / 行标")]),
    ("② 解析引擎", DARK,  GRAY_BG,  [("图层识别","提取墙体、柱、门窗等"),("尺寸标注解析","识别标注线、数值"),("空间语义分析","划定房间、走廊区域"),("文字OCR识别","提取图纸内文字说明")]),
    ("③ 审图引擎", DARK,  RED_BG,   [("防火规范审查","疏散距离、防火分区"),("无障碍审查","坡道、通道尺寸合规"),("建筑强制规定","层高、面积、采光通风"),("结构安全校核","轴网对齐、构件合理性")]),
    ("④ 输出层",   GREEN, RGBColor(0xEC,0xFD,0xF5), [("审图报告 PDF","问题清单+条文依据"),("结构化 JSON","API 接口对接"),("标注图纸","问题位置可视化"),("统计看板","批量数据汇总分析")]),
]
col_w = 2.6
for ci, (hdr, hdr_color, bg_color, items) in enumerate(cols):
    cx = 0.5 + ci * (col_w + 0.7)
    # 箭头（非第一列前）
    if ci > 0:
        ax = cx - 0.58
        add_rect(s2, ax+0.1, 3.5, 0.38, 0.06, hdr_color)
        add_text(s2, "▶", ax+0.32, 3.35, 0.3, 0.35, 18, color=hdr_color)
    add_rect(s2, cx, 0.9, col_w, 0.5, hdr_color)
    add_text(s2, hdr, cx+0.1, 0.94, col_w-0.2, 0.42, 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ii, (it, id_) in enumerate(items):
        iy = 1.55 + ii * 1.32
        add_rect(s2, cx, iy, col_w, 1.15, bg_color, BORDER, 0.5)
        add_text(s2, it, cx+0.12, iy+0.1, col_w-0.24, 0.38, 15, bold=True, color=hdr_color)
        add_text(s2, id_, cx+0.12, iy+0.52, col_w-0.24, 0.55, 12, color=GRAY_TEXT)

# 完成标签
add_rect(s2, 12.1, 3.3, 0.9, 0.9, GREEN)
add_text(s2, "✓", 12.15, 3.32, 0.8, 0.86, 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 幻灯片 3 — 如何进行审图
# ═══════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, 13.33, 7.5, WHITE)
slide_header(s3, "审图流程", BLUE, "如何进行审图")

steps = [
    ("01","上传图纸","支持 DWG/PDF/图片\n拖拽上传，自动格式检测", BLUE),
    ("02","选择规范","按项目类型选择适用规范\n居住/公建/工业/综合体", DARK),
    ("03","自动审查","AI 引擎自动逐条比对\n规范条文，实时生成问题", BLUE),
    ("04","人工复核","审图员对AI结果进行\n确认、修改、补充批注", DARK),
    ("05","报告导出","一键生成正式审图报告\n支持PDF/Word格式", GREEN),
]
sw = 2.4
for si,(num,title,desc,col) in enumerate(steps):
    sx = 0.5 + si*(sw+0.12)
    add_rect(s3, sx, 0.9, sw, 1.85, col)
    add_text(s3, num, sx+0.15, 0.94, sw-0.3, 0.6, 28, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text(s3, title, sx+0.15, 1.55, sw-0.3, 0.4, 18, bold=True, color=WHITE)
    add_text(s3, desc, sx+0.15, 2.0, sw-0.3, 0.7, 13, color=WHITE)

# 下方三大详情卡片
det_data = [
    ("上传界面特性", BLUE, ["最大支持 500MB 单文件","批量上传（最多 50 张）","实时格式校验与预览","自动提取图纸元数据"]),
    ("规范体系覆盖", BLUE, ["GB 50016 建筑防火设计规范","GB 50099 中小学设计规范","GB 50096 住宅设计规范","DB 地方标准（动态更新）"]),
    ("AI 审查能力",  BLUE, ["[核心] 基于 GPT-4V 视觉推理引擎","[扩展] 图规则引擎并行执行","[验证] 置信度评分+人工复核标记",""]),
]
dw = 4.0
for di,(dt,dc,ditems) in enumerate(det_data):
    dx = 0.5 + di*(dw+0.2)
    add_rect(s3, dx, 3.05, dw, 4.15, GRAY_BG, BORDER, 0.5)
    add_text(s3, dt, dx+0.2, 3.15, dw-0.4, 0.4, 16, bold=True, color=dc)
    add_rect(s3, dx+0.2, 3.65, 0.4, 0.05, dc)
    for dii, ditem in enumerate(ditems):
        if ditem:
            add_text(s3, f"● {ditem}", dx+0.2, 3.85+dii*0.75, dw-0.4, 0.65, 13, color=DARK)


# ═══════════════════════════════════════════════
# 幻灯片 4 — 图纸智能解析原理
# ═══════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, 13.33, 7.5, WHITE)
slide_header(s4, "解析技术", DARK, "图纸智能解析原理")

# 左侧深色面板
add_rect(s4, 0.5, 0.9, 6.0, 6.3, DARK)
add_text(s4, "解析管线", 0.8, 1.0, 5.5, 0.5, 22, bold=True, color=WHITE)
add_text(s4, "图纸原文 → 结构化语义数据", 0.8, 1.55, 5.5, 0.35, 14, color=MUTED, italic=True)

pipes = [
    ("STEP 1","几何实体提取","解析线段、圆弧、填充区域形成原始几何图元集合"),
    ("STEP 2","图层语义分类","按图层名规则和ML模型分类墙、柱、门窗、标注等构件"),
    ("STEP 3","拓扑关系构建","推断房间包含关系、门窗归属、走廊连通性拓扑结构"),
    ("STEP 4","尺寸数据提取","解析标注线关联数值，完成比例换算并建立尺寸约束表"),
]
for pi,(badge,ptitle,pdesc) in enumerate(pipes):
    py = 2.05 + pi*1.1
    bg = GREEN if pi==3 else BLUE
    add_rect(s4, 0.7, py, 5.6, 0.95, RGBColor(0x26,0x26,0x26))
    add_rect(s4, 0.85, py+0.1, 0.75, 0.3, bg)
    add_text(s4, badge, 0.87, py+0.11, 0.71, 0.28, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, ptitle, 1.68, py+0.08, 4.4, 0.34, 15, bold=True, color=WHITE)
    add_text(s4, pdesc, 0.85, py+0.46, 5.3, 0.4, 12, color=MUTED)

# 右侧
add_rect(s4, 6.85, 0.9, 6.0, 2.5, GRAY_BG)
add_text(s4, "核心技术栈", 7.1, 1.0, 5.5, 0.42, 18, bold=True, color=DARK)
add_rect(s4, 7.1, 1.5, 0.4, 0.06, BLUE)
techs = [("ezdxf","DWG/DXF解析"),("PyMuPDF","PDF矢量提取"),("YOLO v8","构件目标检测"),("GPT-4V","语义推理审查")]
for ti,(tn,td) in enumerate(techs):
    tx = 6.95 + (ti%2)*2.9
    ty = 1.65 + (ti//2)*0.95
    add_rect(s4, tx, ty, 2.65, 0.8, WHITE, BORDER, 0.5)
    add_text(s4, tn, tx+0.12, ty+0.06, 2.4, 0.35, 15, bold=True, color=BLUE)
    add_text(s4, td, tx+0.12, ty+0.42, 2.4, 0.3, 12, color=GRAY_TEXT)

add_rect(s4, 6.85, 3.55, 6.0, 1.9, DARK)
add_text(s4, "解析结果示例（JSON）", 7.1, 3.65, 5.5, 0.38, 15, bold=True, color=WHITE)
code_lines = ['{ "room": "走廊",','  "width_mm": 1500,','  "area_m2": 24.6,','  "fire_exit": true,','  "compliance": "PASS"','}']
for li,ln in enumerate(code_lines):
    c = RGBColor(0x86,0xEF,0xAC) if "PASS" in ln else (WHITE if ln.strip().startswith('"') else MUTED)
    add_text(s4, ln, 7.1, 4.1+li*0.22, 5.5, 0.22, 12, color=c, italic=("PASS" not in ln and li>0))

add_rect(s4, 6.85, 5.6, 6.0, 1.6, BLUE_BG)
acc_data = [("99.2%","线段识别率"),("96.8%","标注解析率"),("94.1%","语义分类精度")]
for ai,(an,al) in enumerate(acc_data):
    ax = 7.0 + ai*2.0
    add_text(s4, an, ax, 5.75, 1.8, 0.6, 28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s4, al, ax, 6.38, 1.8, 0.35, 13, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 幻灯片 5 — 最终审查成果展示
# ═══════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, 13.33, 7.5, WHITE)
slide_header(s5, "成功结果", GREEN, "最终审查成果展示")

# KPI 条
kpis = [("1,280","累计审图张数",BLUE,WHITE),("98.3%","问题识别准确率",DARK,WHITE),("72%","人力审图时间节省",GREEN,WHITE),("24s","平均单图处理时长",GRAY_BG,BLUE)]
for ki,(kn,kl,kbg,ktc) in enumerate(kpis):
    kx = 0.5 + ki*3.25
    add_rect(s5, kx, 0.9, 3.0, 1.4, kbg)
    add_text(s5, kn, kx+0.1, 0.96, 2.8, 0.82, 38, bold=True, color=ktc, align=PP_ALIGN.CENTER)
    add_text(s5, kl, kx+0.1, 1.82, 2.8, 0.4, 14, color=GRAY_TEXT if kbg==GRAY_BG else WHITE, align=PP_ALIGN.CENTER)

# 案例1
cx1 = 0.5
add_rect(s5, cx1, 2.55, 6.0, 4.7, GRAY_BG, BORDER, 0.5)
add_rect(s5, cx1+0.2, 2.7, 1.4, 0.35, GREEN_BG)
add_text(s5, "✓ 审查通过", cx1+0.25, 2.72, 1.3, 0.3, 12, bold=True, color=GREEN)
add_text(s5, "某住宅小区 A 座", cx1+0.2, 3.15, 5.6, 0.45, 20, bold=True, color=DARK)
add_text(s5, "地上 18 层 / 建筑面积 12,600㎡", cx1+0.2, 3.65, 5.6, 0.35, 13, color=GRAY_TEXT)
add_rect(s5, cx1+0.2, 4.1, 5.6, 0.02, BORDER)
rows1 = [("✓","防火规范：走廊宽度 1800mm ≥ 1200mm",GREEN),("✓","疏散距离：最远点 32m ≤ 40m",GREEN),("✓","无障碍：坡道坡度 1:12 合规",GREEN),("✓","层高：标准层 2950mm ≥ 2800mm",GREEN)]
for ri,(ic,rt,rc) in enumerate(rows1):
    add_text(s5, f"{ic}  {rt}", cx1+0.2, 4.2+ri*0.42, 5.6, 0.38, 13, color=rc)
sumdata1 = [("47","检查项",GREEN_BG,GREEN),("47","通过",GREEN_BG,GREEN),("0","问题项",GRAY_BG,MUTED)]
for si2,(sn,sl,sbg,sc) in enumerate(sumdata1):
    sx2 = cx1+0.2+si2*1.9
    add_rect(s5, sx2, 6.4, 1.6, 0.7, sbg)
    add_text(s5, sn, sx2+0.1, 6.44, 1.4, 0.38, 22, bold=True, color=sc, align=PP_ALIGN.CENTER)
    add_text(s5, sl, sx2+0.1, 6.82, 1.4, 0.25, 12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# 案例2
cx2 = 6.85
add_rect(s5, cx2, 2.55, 6.0, 4.7, GRAY_BG, BORDER, 0.5)
add_rect(s5, cx2+0.2, 2.7, 1.7, 0.35, ORANGE_BG)
add_text(s5, "⚠ 有条件通过", cx2+0.25, 2.72, 1.6, 0.3, 12, bold=True, color=ORANGE)
add_text(s5, "某公共建筑 B 项目", cx2+0.2, 3.15, 5.6, 0.45, 20, bold=True, color=DARK)
add_text(s5, "地上 5 层商业综合体 / 建筑面积 8,200㎡", cx2+0.2, 3.65, 5.6, 0.35, 13, color=GRAY_TEXT)
add_rect(s5, cx2+0.2, 4.1, 5.6, 0.02, BORDER)
rows2 = [("✓","防火分区面积：符合规范要求",GREEN),("⚠","无障碍卫生间：1400mm 不足 1500mm",ORANGE),("⚠","疏散楼梯B区：净宽 1050mm < 1100mm",ORANGE),("✓","采光面积比：主要房间均 ≥ 1/7",GREEN)]
for ri,(ic,rt,rc) in enumerate(rows2):
    add_text(s5, f"{ic}  {rt}", cx2+0.2, 4.2+ri*0.42, 5.6, 0.38, 13, color=rc)
sumdata2 = [("52","检查项",GREEN_BG,GREEN),("50","通过",GREEN_BG,GREEN),("2","待修改",RGBColor(0xFE,0xF9,0xC3),ORANGE)]
for si2,(sn,sl,sbg,sc) in enumerate(sumdata2):
    sx2 = cx2+0.2+si2*1.9
    add_rect(s5, sx2, 6.4, 1.6, 0.7, sbg)
    add_text(s5, sn, sx2+0.1, 6.44, 1.4, 0.38, 22, bold=True, color=sc, align=PP_ALIGN.CENTER)
    add_text(s5, sl, sx2+0.1, 6.82, 1.4, 0.25, 12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# 幻灯片 6 — 生成的审图报告样例
# ═══════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_rect(s6, 0, 0, 13.33, 7.5, WHITE)
slide_header(s6, "审图报告", BLUE, "生成的审图报告样例")

# 左侧报告样本
add_rect(s6, 0.5, 0.9, 5.5, 6.4, WHITE, BORDER, 0.8)
# 报告标题栏
add_rect(s6, 0.5, 0.9, 5.5, 0.75, BLUE)
add_text(s6, "建筑施工图审查报告", 0.7, 0.96, 3.8, 0.38, 16, bold=True, color=WHITE)
add_text(s6, "Report No. SH-2025-0312", 0.7, 1.38, 3.8, 0.24, 11, color=RGBColor(0xBF,0xDB,0xFF))
add_rect(s6, 4.6, 1.0, 1.0, 0.45, WHITE)
add_text(s6, "已签发", 4.65, 1.04, 0.9, 0.37, 13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
# 项目信息栏
add_rect(s6, 0.5, 1.65, 5.5, 0.65, GRAY_BG)
add_text(s6, "项目名称", 0.65, 1.68, 1.7, 0.25, 11, color=MUTED)
add_text(s6, "某住宅小区 A 座", 0.65, 1.93, 1.7, 0.3, 13, bold=True, color=DARK)
add_text(s6, "审查日期", 2.5, 1.68, 1.5, 0.25, 11, color=MUTED)
add_text(s6, "2025-03-12", 2.5, 1.93, 1.5, 0.3, 13, bold=True, color=DARK)
add_text(s6, "结论", 4.2, 1.68, 1.5, 0.25, 11, color=MUTED)
add_text(s6, "符合规范", 4.2, 1.93, 1.5, 0.3, 14, bold=True, color=GREEN)

# 审查详情
add_rect(s6, 0.5, 2.3, 5.5, 0.02, BORDER)
add_text(s6, "一、防火规范检查", 0.65, 2.38, 5.2, 0.35, 14, bold=True, color=DARK)
chk1 = [("通过","疏散走廊净宽 ≥ 1200mm","实测 1800mm"),("通过","疏散距离 ≤ 40m（一类高层）","最远 32m"),("通过","防火门设置（乙级）","已配置")]
for ci,(cs,ct,cv) in enumerate(chk1):
    cy = 2.82+ci*0.42
    add_rect(s6, 0.65, cy, 0.6, 0.3, GREEN_BG)
    add_text(s6, cs, 0.66, cy+0.02, 0.58, 0.26, 11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s6, ct, 1.32, cy+0.02, 3.0, 0.3, 13, color=DARK)
    add_text(s6, cv, 4.4, cy+0.02, 1.4, 0.3, 12, color=GRAY_TEXT)

add_rect(s6, 0.5, 4.12, 5.5, 0.02, BORDER)
add_text(s6, "二、无障碍设施检查", 0.65, 4.2, 5.2, 0.35, 14, bold=True, color=DARK)
chk2 = [("通过","无障碍坡道坡度 ≤ 1:12","实测 1:14"),("通过","无障碍停车位数量","3 个（≥ 2）")]
for ci,(cs,ct,cv) in enumerate(chk2):
    cy = 4.65+ci*0.42
    add_rect(s6, 0.65, cy, 0.6, 0.3, GREEN_BG)
    add_text(s6, cs, 0.66, cy+0.02, 0.58, 0.26, 11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s6, ct, 1.32, cy+0.02, 3.0, 0.3, 13, color=DARK)
    add_text(s6, cv, 4.4, cy+0.02, 1.4, 0.3, 12, color=GRAY_TEXT)

# 结论栏
add_rect(s6, 0.5, 6.55, 5.5, 0.75, GRAY_BG)
add_rect(s6, 0.65, 6.68, 0.9, 0.32, GREEN_BG)
add_text(s6, "✓ 审查完成", 0.66, 6.69, 0.88, 0.3, 11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s6, "全部 47 项均符合相关规范要求，建议予以通过。", 1.65, 6.68, 4.2, 0.4, 13, color=GRAY_TEXT)

# 右侧报告说明
add_rect(s6, 6.3, 0.9, 6.5, 3.0, GRAY_BG)
add_text(s6, "报告结构", 6.55, 1.0, 6.0, 0.42, 18, bold=True, color=DARK)
add_rect(s6, 6.55, 1.52, 0.4, 0.06, BLUE)
chapters = [("01","项目基本信息"),("02","审查依据（规范清单）"),("03","逐条审查详情"),("04","问题汇总与建议"),("05","审图结论与签章")]
for chi,(cn,ct) in enumerate(chapters):
    cy2 = 1.65+chi*0.42
    add_rect(s6, 6.55, cy2, 5.9, 0.36, WHITE, BORDER, 0.5)
    add_text(s6, cn, 6.68, cy2+0.04, 0.45, 0.28, 12, bold=True, color=BLUE)
    add_text(s6, ct, 7.25, cy2+0.04, 5.0, 0.28, 13, color=DARK)

add_rect(s6, 6.3, 4.05, 6.5, 1.55, BLUE_BG)
add_text(s6, "导出格式", 6.55, 4.15, 6.0, 0.38, 16, bold=True, color=BLUE)
fmts = [("PDF","正式报告"),("DOCX","可编辑版"),("JSON","API对接")]
for fi,(fn,fl) in enumerate(fmts):
    fx = 6.45+fi*2.1
    add_rect(s6, fx, 4.65, 1.85, 0.8, WHITE)
    add_text(s6, fn, fx+0.1, 4.72, 1.65, 0.4, 20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s6, fl, fx+0.1, 5.12, 1.65, 0.28, 12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

add_rect(s6, 6.3, 5.75, 6.5, 1.55, RGBColor(0xEC,0xFD,0xF5))
add_text(s6, "本次审查亮点", 6.55, 5.85, 6.0, 0.38, 16, bold=True, color=RGBColor(0x06,0x5F,0x46))
highlights = ["✓  全流程自动化，无需手动填写报告模板","✓  问题定位精确到坐标，方便设计师修改","✓  条文依据一一对应，审批流程可追溯"]
for hi,ht in enumerate(highlights):
    add_text(s6, ht, 6.55, 6.28+hi*0.32, 5.9, 0.3, 13, color=DARK)


# ═══════════════════════════════════════════════
# 幻灯片 7 — 结束页
# ═══════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_rect(s7, 0, 0, 13.33, 7.5, DARK)
add_rect(s7, 0, 0, 0.08, 7.5, BLUE)
# 渐变装饰圆
add_rect(s7, 8.5, -1.0, 5.5, 5.5, RGBColor(0x06,0x2D,0x6B))

# 中心内容
add_rect(s7, 3.5, 1.8, 4.0, 0.65, BLUE)
add_text(s7, "智能审图 · 赋能设计", 3.65, 1.86, 3.7, 0.53, 22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s7, "感谢观看", 2.5, 2.65, 8.33, 1.6, 72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s7, "让每一张图纸都经过严格审查，让每一栋建筑都符合安全规范",
         2.0, 4.4, 9.33, 0.65, 20, color=MUTED, align=PP_ALIGN.CENTER)
add_rect(s7, 6.05, 5.15, 1.23, 0.06, BLUE)
add_text(s7, "📧 review@sparkflow.ai", 3.2, 5.4, 3.2, 0.38, 16, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s7, "🌐 www.sparkflow.ai",    6.9, 5.4, 3.0, 0.38, 16, color=MUTED, align=PP_ALIGN.CENTER)

# 页脚
add_text(s7, "SparkFlow 智能审图平台", 0.5, 7.05, 6.0, 0.3, 13, color=RGBColor(0x4B,0x55,0x63))
add_text(s7, "演示汇报文档 · 2025 · 保密", 7.0, 7.05, 6.0, 0.3, 13, color=RGBColor(0x4B,0x55,0x63), align=PP_ALIGN.RIGHT)

# ─────────────────── 保存 ───────────────────
output = r"D:\code\project\moment\SparkFlow\docs\智能审图汇报演示.pptx"
prs.save(output)
print(f"✅ PPT 已保存至：{output}")
